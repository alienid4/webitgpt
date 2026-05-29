from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(r"F:\ClaudeHome\webitgpt\docs\presentations")
OUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX = OUT_DIR / "it_war_room_apple_keynote_v3.pptx"
NOTES = OUT_DIR / "it_war_room_apple_keynote_v3_talk_track.md"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLORS = {
    "bg": RGBColor(250, 249, 245),
    "ink": RGBColor(20, 20, 19),
    "muted": RGBColor(112, 112, 112),
    "teal": RGBColor(38, 168, 137),
    "teal_dark": RGBColor(13, 92, 71),
    "orange": RGBColor(232, 124, 7),
    "red": RGBColor(224, 11, 0),
    "line": RGBColor(190, 190, 190),
    "soft": RGBColor(235, 247, 232),
    "white": RGBColor(255, 255, 255),
    "blue": RGBColor(33, 119, 255),
}
FONT = "Microsoft JhengHei"
MONO = "Consolas"


def blank():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    return slide


def set_text(tf, text, size=28, color="ink", bold=False, font=FONT, align=None):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLORS[color]
    if align:
        p.alignment = align
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = COLORS[color]
    return p


def textbox(slide, x, y, w, h, text, size=28, color="ink", bold=False, align=None, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    set_text(tf, text, size, color, bold, font, align)
    return box


def title(slide, text, sub=None):
    textbox(slide, 0.65, 0.45, 12, 0.7, text, 28, "ink", True)
    if sub:
        textbox(slide, 0.67, 1.15, 11, 0.4, sub, 12.5, "muted")


def pill(slide, x, y, w, h, text, border="teal", fill="white", color="teal_dark", size=15, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS[fill]
    shp.line.color.rgb = COLORS[border]
    shp.line.width = Pt(1.4)
    shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    shp.text_frame.margin_left = Inches(0.08)
    shp.text_frame.margin_right = Inches(0.08)
    set_text(shp.text_frame, text, size, color, bold, align=PP_ALIGN.CENTER)
    return shp


def card(slide, x, y, w, h, head, body="", accent="teal", value=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS["white"]
    shp.line.color.rgb = COLORS["line"]
    shp.line.width = Pt(0.9)
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS[accent]
        bar.line.fill.background()
    textbox(slide, x + 0.22, y + 0.18, w - 0.35, 0.28, head, 12, "teal_dark", True)
    if value is not None:
        textbox(slide, x + 0.22, y + 0.58, w - 0.35, 0.55, str(value), 32, "ink", True, font=MONO)
        textbox(slide, x + 0.22, y + 1.22, w - 0.35, h - 1.35, body, 11.5, "muted")
    else:
        textbox(slide, x + 0.22, y + 0.62, w - 0.35, h - 0.8, body, 15, "ink")
    return shp


def line(slide, x1, y1, x2, y2, color="teal_dark", width=2.2, dash=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = COLORS[color]
    conn.line.width = Pt(width)
    if dash:
        conn.line.dash_style = 4
    return conn


def circle(slide, x, y, d, text, border="teal", fill="soft", color="ink", size=18, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS[fill]
    shp.line.color.rgb = COLORS[border]
    shp.line.width = Pt(2)
    shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(shp.text_frame, text, size, color, bold, align=PP_ALIGN.CENTER)
    return shp


def footer(slide, n):
    textbox(slide, 0.65, 7.08, 8, 0.22, f"資訊戰情室 / Apple-style concept v3 / {n:02d}", 8.5, "muted")


slide = blank()
textbox(slide, 0.7, 2.05, 11.8, 0.9, "資訊戰情室", 52, "ink", True)
textbox(slide, 0.75, 3.05, 10.8, 0.55, "一個畫面，看見資產、健康與風險", 25, "teal_dark")
pill(slide, 0.75, 4.05, 2.2, 0.45, "IT Inspection Command Center", "teal", "soft", "teal_dark", 11)
footer(slide, 1)

slide = blank()
title(slide, "今天的問題，不是沒有資料。", "是資料分散、格式不同，而且很難相信。")
for i, (head, body) in enumerate(
    [("資產", "Excel、人工填寫、IP 衝突"), ("巡檢", "每天看很多畫面，卻缺少結論"), ("帳號", "高權限、PAM、異動難追"), ("故障", "出事後才開始找影響範圍")]
):
    card(slide, 0.85 + i * 3.05, 2.35, 2.65, 2.2, head, body, "orange")
textbox(slide, 1.0, 5.65, 11.2, 0.55, "結果：維運人員花時間整理資料，主管仍然看不到真正風險。", 24, "ink", True, PP_ALIGN.CENTER)
footer(slide, 2)

slide = blank()
textbox(slide, 1.0, 1.55, 11.2, 0.8, "把 IT 狀態，收斂成一個可信任畫面。", 39, "ink", True, PP_ALIGN.CENTER)
for i, label in enumerate(["資產", "開門檢查", "帳號盤點", "效能月報", "系統拓撲"]):
    circle(slide, 1.05 + i * 2.45, 3.3, 1.35, label, "teal", "white", "teal_dark", 15)
    if i < 4:
        line(slide, 2.42 + i * 2.45, 3.98, 3.35 + i * 2.45, 3.98, "teal", 2.2)
textbox(slide, 1.3, 5.6, 10.7, 0.4, "不是取代工程師，而是讓工程師和主管看同一份證據。", 22, "muted", False, PP_ALIGN.CENTER)
footer(slide, 3)

slide = blank()
title(slide, "每天開門第一眼，知道哪裡要處理。")
for x, head, value, body, accent in [
    (0.8, "總體狀態", "100", "正常 8 / 警示 1 / 待檢 0", "teal"),
    (4.05, "磁碟警戒", "77%", "最高 filesystem：/ 77%", "orange"),
    (7.3, "帳號風險", "1", "服務帳號可登入，需複核", "orange"),
    (9.9, "本日巡檢", "3", "已完成 3 台", "teal"),
]:
    card(slide, x, 1.75, 3.0 if x < 7 else 2.55, 1.25, head, body, accent, value)
for row in range(3):
    for col in range(3):
        idx = row * 3 + col
        labels = ["連線", "CPU / MEM / SWAP", "檔案系統", "程序", "服務", "帳號", "安全設定", "套件", "系統日誌"]
        card(slide, 0.95 + col * 4.05, 3.5 + row * 0.78, 3.35, 0.58, labels[idx], "OK" if idx != 8 else "WARN", "teal" if idx != 8 else "orange")
footer(slide, 4)

slide = blank()
title(slide, "資產管理，先解決「可信任」。", "人工資料、掃描結果、實際盤點，必須能對帳。")
for i, (head, body) in enumerate(
    [("人工匯入", "CMDB / Excel / CSV"), ("主動掃描", "IPAM / nmap 建草稿"), ("實際盤點", "SSH / WinRM / AIX SSH"), ("治理狀態", "等待防火牆 / 弱掃 / PAM")]
):
    card(slide, 0.9 + i * 3.0, 2.1, 2.55, 1.55, head, body, "teal")
    if i < 3:
        line(slide, 3.48 + i * 3.0, 2.85, 3.9 + i * 3.0, 2.85, "teal", 2)
textbox(slide, 1.0, 5.0, 11.2, 0.7, "有 IP 但沒有資產、資產下線但未註記，會變成清單，而不是藏在 Excel 裡。", 28, "ink", True, PP_ALIGN.CENTER)
footer(slide, 5)

slide = blank()
title(slide, "開門檢查，只講重點。", "每台主機先看小卡；需要時再展開證據。")
for i, (name, health, warn) in enumerate([("巡檢系統主機", "100", "0"), ("受監控主機-Rocky", "86", "1"), ("受監控主機-Debian", "92", "1")]):
    x = 0.8 + i * 4.15
    card(slide, x, 1.75, 3.55, 3.95, "HOST DASHBOARD", "", "teal" if warn == "0" else "orange")
    textbox(slide, x + 0.25, 2.15, 2.4, 0.35, name, 16, "ink", True)
    circle(slide, x + 0.25, 2.75, 1.15, health, "teal" if warn == "0" else "orange", "white", "ink", 18)
    textbox(slide, x + 1.65, 2.85, 1.4, 0.3, "正常" if warn == "0" else "警示", 16, "teal_dark" if warn == "0" else "orange", True)
    textbox(slide, x + 0.25, 4.15, 3.0, 0.6, "CPU 9%\nMemory 74%\nFilesystem / 77%", 13, "ink")
    pill(slide, x + 0.25, 5.35, 1.35, 0.38, "執行 L3", "teal", "teal", "white", 11)
footer(slide, 6)

slide = blank()
title(slide, "深度檢查，不是黑箱。", "PASS 也要有證據；WARN 要說清楚問題、證據與處理方式。")
card(slide, 0.9, 1.8, 3.25, 3.75, "問題點", "網卡 dropped counter 增加，可能影響連線品質。", "orange")
card(slide, 4.45, 1.8, 3.25, 3.75, "證據", "NET-02：dropped=17\nFilesystem：/ 77%\nService：無 failed service", "teal")
card(slide, 8.0, 1.8, 4.1, 3.75, "建議處置", "1. 確認哪張網卡增加\n2. 若持續增加，請 VM / 網路管理員檢查\n3. 修正後重跑深度檢查", "teal")
footer(slide, 7)

slide = blank()
title(slide, "帳號盤點，要回答主管真正想問的。")
for i, (value, label) in enumerate([("93", "帳號總數"), ("12", "高權限"), ("11", "PAM 納管"), ("1", "服務帳號可登入"), ("0", "新增異動")]):
    card(slide, 0.75 + i * 2.45, 1.8, 2.1, 1.25, label, "", "teal" if i != 3 else "orange", value)
textbox(slide, 1.0, 4.1, 11.2, 0.55, "重點不是列出所有帳號，而是知道誰需要複核、誰已納管、誰和上次不同。", 27, "ink", True, PP_ALIGN.CENTER)
pill(slide, 3.6, 5.35, 2.0, 0.48, "差異報告", "teal", "soft", "teal_dark", 14)
pill(slide, 5.8, 5.35, 2.0, 0.48, "PAM 狀態", "teal", "soft", "teal_dark", 14)
pill(slide, 8.0, 5.35, 2.0, 0.48, "高權限清冊", "teal", "soft", "teal_dark", 14)
footer(slide, 8)

slide = blank()
title(slide, "效能月報，讓主管看到趨勢，不看 raw data。")
for x, head, value, body, accent in [
    (0.8, "納管主機", "3", "本月有 NMON 資料", "teal"),
    (3.65, "採樣覆蓋", "96%", "平均覆蓋率", "teal"),
    (6.5, "容量風險", "1", "需追蹤主機", "orange"),
]:
    card(slide, x, 1.7, 2.6, 1.25, head, body, accent, value)
x0, y0, width, height = 1.0, 3.55, 10.8, 2.1
line(slide, x0, y0 + height, x0 + width, y0 + height, "line", 1)
line(slide, x0, y0, x0, y0 + height, "line", 1)
points = [(x0 + 0.3, y0 + 1.6), (x0 + 2.1, y0 + 1.4), (x0 + 3.8, y0 + 1.2), (x0 + 5.6, y0 + 1.05), (x0 + 7.5, y0 + 0.82), (x0 + 10.0, y0 + 0.72)]
for a, b in zip(points, points[1:]):
    line(slide, a[0], a[1], b[0], b[1], "teal", 3)
for point in points:
    circle(slide, point[0] - 0.06, point[1] - 0.06, 0.12, "", "teal", "teal", "white", 1)
textbox(slide, 1.0, 6.0, 10.5, 0.35, "呈現 CPU、記憶體、磁碟、尖峰與容量趨勢；raw file 留在明細，不放在主管頁。", 17, "muted")
footer(slide, 9)

slide = blank()
title(slide, "系統拓撲，回答「誰會被影響」。")
circle(slide, 0.95, 2.95, 1.35, "核心\n系統", "red", "white", "ink", 16)
circle(slide, 3.45, 1.65, 1.25, "交易\n核心", "blue", "white", "ink", 14)
circle(slide, 3.45, 3.2, 1.25, "巡檢\n系統", "orange", "white", "ink", 14)
circle(slide, 3.45, 4.75, 1.25, "帳務\n核心", "blue", "white", "ink", 14)
for yy in [2.25, 3.8, 5.35]:
    line(slide, 2.3, 3.63, 3.45, yy, "teal", 2)
for i, (name, yy) in enumerate([("AP", 1.45), ("DB", 2.45), ("Batch", 3.45), ("FE", 4.45), ("PAM", 5.45)]):
    pill(slide, 6.15, yy, 1.75, 0.55, name, "teal", "white", "teal_dark", 14)
    line(slide, 4.7, 3.8, 6.15, yy + 0.28, "teal", 2 if i < 3 else 1.2, dash=i >= 3)
for i, ip in enumerate(["10.1.4.21", "10.1.4.22", "10.1.4.23"]):
    pill(slide, 9.25, 2.0 + i * 0.9, 1.95, 0.5, ip, "teal", "soft", "ink", 12)
    line(slide, 7.9, 2.0 + i * 0.9 + 0.25, 9.25, 2.0 + i * 0.9 + 0.25, "teal", 1.7)
textbox(slide, 8.9, 5.45, 3.7, 0.55, "維護前先看影響範圍，故障時知道要通知誰。", 21, "ink", True)
footer(slide, 10)

slide = blank()
title(slide, "AI 可以協助 debug，但不直接碰 VM。")
for i, (head, body) in enumerate([("VM 產生", "log / report"), ("去識別化", "bundle"), ("公司 GPT", "分析原因"), ("Codex", "修 code"), ("測試機", "重新驗證")]):
    card(slide, 0.65 + i * 2.5, 2.45, 2.05, 1.45, head, body, "teal")
    if i < 4:
        line(slide, 2.72 + i * 2.5, 3.15, 3.12 + i * 2.5, 3.15, "teal", 2.3)
textbox(slide, 1.1, 5.3, 11.0, 0.5, "原則：AI 看整理後的證據，不直接登入公司 VM。", 29, "ink", True, PP_ALIGN.CENTER)
footer(slide, 11)

slide = blank()
title(slide, "架構很簡單：收集、判斷、呈現、追蹤。")
for i, (head, body) in enumerate([("收集", "SSH / WinRM / nmap / NMON"), ("判斷", "規則 / 差異 / 門檻"), ("呈現", "Dashboard / 報表 / 拓撲"), ("追蹤", "MongoDB / Log / Audit")]):
    card(slide, 1.0 + i * 3.0, 2.5, 2.45, 1.55, head, body, "teal")
    if i < 3:
        line(slide, 3.48 + i * 3.0, 3.25, 3.95 + i * 3.0, 3.25, "teal", 2.3)
textbox(slide, 1.1, 5.4, 11.0, 0.55, "技術細節可以拆開，但主管只需要知道：資料從哪裡來、怎麼判斷、結果在哪裡看。", 22, "muted", False, PP_ALIGN.CENTER)
footer(slide, 12)

slide = blank()
title(slide, "落地規劃：先把每天會用的做穩。")
for i, (phase, body) in enumerate([("Phase 1", "資產、帳號、開門檢查"), ("Phase 2", "深度檢查、效能月報"), ("Phase 3", "拓撲、通知、AI Debug"), ("Phase 4", "正式上線、權限與稽核")]):
    card(slide, 0.9 + i * 3.05, 2.2, 2.6, 2.2, phase, body, "teal" if i < 2 else "orange")
textbox(slide, 1.0, 5.45, 11.2, 0.6, "交付對象：維運人員每天使用；主管與稽核看報表與風險清單。", 26, "ink", True, PP_ALIGN.CENTER)
footer(slide, 13)

slide = blank()
textbox(slide, 1.0, 2.25, 11.3, 0.8, "讓 IT 從被動救火，變成主動掌握。", 42, "ink", True, PP_ALIGN.CENTER)
textbox(slide, 1.6, 3.55, 10.0, 0.45, "資訊戰情室", 26, "teal_dark", True, PP_ALIGN.CENTER)
footer(slide, 14)

prs.save(PPTX)

NOTES.write_text(
    """# 資訊戰情室 Apple-style PPT 講稿 v3

1. 封面：資訊戰情室不是另一個後台，而是一個 IT 狀態入口。
2. 痛點：問題不是沒有資料，是資料分散且難以相信。
3. 解法：把資產、巡檢、帳號、效能、拓撲收斂到同一畫面。
4. 開門第一眼：每天先知道哪裡需要處理。
5. 資產：處理人工資料與實際狀態不一致。
6. 開門檢查：小卡先講重點，明細留給工程師。
7. 深度檢查：PASS/WARN 都要有證據。
8. 帳號盤點：看高權限、PAM、差異和需要複核的帳號。
9. 效能月報：主管看趨勢，raw data 留明細。
10. 拓撲：維護或故障時知道影響誰。
11. AI Debug Loop：AI 協助分析，但不直接碰 VM。
12. 架構：收集、判斷、呈現、追蹤。
13. 落地：先把每天會用的做穩。
14. 收尾：從被動救火，變成主動掌握。
""",
    encoding="utf-8",
)

print(PPTX)
